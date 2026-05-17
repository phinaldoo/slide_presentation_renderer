from __future__ import annotations

from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path
from urllib.parse import urlsplit
import uuid

from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches

# Configuration
SELECTOR = ".slide"  # CSS selector for 16:9 slide sections, e.g. "section" or ".slide"
SLIDE_BLEED = Inches(0.02)  # Slight overlap to avoid visible slide edges


def _is_allowed_request_url(request_url: str, allowed_origin: str | None) -> bool:
    """Check if request URL is allowed based on origin policy."""
    if request_url.startswith(("data:", "blob:", "about:")):
        return True
    if not allowed_origin:
        return True
    return _extract_origin(request_url) == _extract_origin(allowed_origin)


def _extract_origin(url: str) -> str | None:
    """Extract normalized origin from URL."""
    try:
        parsed = urlsplit(url)
    except ValueError:
        return None

    if not parsed.scheme or not parsed.netloc:
        return None

    return f"{parsed.scheme.lower()}://{parsed.netloc.lower()}"


async def _apply_request_guard(context, allowed_origin: str | None) -> None:
    """Apply request guard to block disallowed origins."""
    if not allowed_origin:
        return

    async def _route_handler(route) -> None:
        request_url = route.request.url
        if _is_allowed_request_url(request_url, allowed_origin):
            await route.continue_()
            return
        await route.abort()

    await context.route("**/*", _route_handler)


def _create_presentation() -> Presentation:
    """Create a new 16:9 PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    return prs


def _save_presentation(prs: Presentation, output_dir: str | Path | None) -> Path:
    """Save presentation to output directory with generated filename."""
    default_dir = Path("presentations") / datetime.now(timezone.utc).strftime("%Y-%m-%d")
    output_path = Path(output_dir) if output_dir else default_dir
    output_path.mkdir(parents=True, exist_ok=True)
    file_path = output_path / f"presentation_{uuid.uuid4()}.pptx"
    prs.save(file_path)
    return file_path


async def _render_page_to_pptx(
    page,
    prs: Presentation,
    selector_to_use: str,
    max_slides: int | None = None,
) -> None:
    """Render page elements matching selector to PPTX slides."""
    sections = await page.locator(selector_to_use).all()
    count = len(sections)
    if count == 0:
        raise RuntimeError(
            f"No elements were found for selector '{selector_to_use}'. "
            "Does the selector match your HTML file?"
        )
    if max_slides is not None and count > max_slides:
        raise ValueError(f"too many slides rendered (max {max_slides})")

    for section in sections:
        await section.scroll_into_view_if_needed()
        image_bytes = await section.screenshot()

        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        image_stream = BytesIO(image_bytes)
        slide.shapes.add_picture(
            image_stream,
            left=-SLIDE_BLEED,
            top=-SLIDE_BLEED,
            width=prs.slide_width + SLIDE_BLEED * 2,
            height=prs.slide_height + SLIDE_BLEED * 2,
        )


async def html_to_pptx(
    html_content: str,
    *,
    output_dir: str | Path | None = None,
    selector: str | None = None,
    allowed_origin: str | None = None,
    page_load_timeout_ms: int = 30_000,
    max_slides: int | None = None,
) -> Path:
    """Render HTML slides to a PPTX file and return the saved file path."""

    selector_to_use = selector or SELECTOR

    prs = _create_presentation()

    async with async_playwright() as p:
        browser = await p.chromium.launch()
        context = None
        try:
            context = await browser.new_context(viewport={"width": 1920, "height": 1080})
            await _apply_request_guard(context, allowed_origin)
            page = await context.new_page()
            page.set_default_navigation_timeout(page_load_timeout_ms)
            page.set_default_timeout(page_load_timeout_ms)

            if html_content and html_content.strip():
                await page.set_content(
                    html_content,
                    wait_until="networkidle",
                    timeout=page_load_timeout_ms,
                )

            await page.wait_for_load_state("networkidle", timeout=page_load_timeout_ms)
            await _render_page_to_pptx(page, prs, selector_to_use, max_slides)
        finally:
            if context is not None:
                await context.close()
            await browser.close()

    return _save_presentation(prs, output_dir)


async def html_url_to_pptx(
    html_url: str,
    *,
    output_dir: str | Path | None = None,
    selector: str | None = None,
    allowed_origin: str | None = None,
    page_load_timeout_ms: int = 30_000,
    max_slides: int | None = None,
) -> Path:
    """Load HTML from a URL and convert it to PPTX."""

    selector_to_use = selector or SELECTOR
    prs = _create_presentation()

    async with async_playwright() as p:
        browser = await p.chromium.launch()
        context = None
        try:
            context = await browser.new_context(viewport={"width": 1920, "height": 1080})
            await _apply_request_guard(context, allowed_origin)
            page = await context.new_page()
            page.set_default_navigation_timeout(page_load_timeout_ms)
            page.set_default_timeout(page_load_timeout_ms)
            await page.goto(
                html_url,
                wait_until="networkidle",
                timeout=page_load_timeout_ms,
            )
            await page.wait_for_load_state("networkidle", timeout=page_load_timeout_ms)
            await _render_page_to_pptx(page, prs, selector_to_use, max_slides)
        finally:
            if context is not None:
                await context.close()
            await browser.close()

    return _save_presentation(prs, output_dir)


async def html_file_to_pptx(
    html_file_path: str | Path,
    *,
    output_dir: str | Path | None = None,
    selector: str | None = None,
) -> Path:
    """Read an HTML file and convert it to PPTX."""

    html_path = Path(html_file_path)
    if not html_path.exists():
        raise FileNotFoundError(f"HTML file '{html_path}' was not found.")

    html_content = html_path.read_text(encoding="utf-8")
    return await html_to_pptx(
        html_content,
        output_dir=output_dir,
        selector=selector,
    )
